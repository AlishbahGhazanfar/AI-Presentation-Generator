import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from fastapi import FastAPI, Form, HTTPException
from fastapi.responses import FileResponse
import speech_recognition as sr
from pptx.enum.text import PP_ALIGN



# Set up OpenAI API key
openai.api_key = "sk-svcacct-MklMoHAlVpjrQXnU8j_0FM_S3UpIG3NZqUnGacRDVjO-9d6CIGVIhpMMtjO800AT3BlbkFJRD2guGV26y93k0zPsKWPdfXrG2EEcOwvwlWZqZdzjGSklQ454mXu5t4VFg7aJAA"  # Replace with your OpenAI API key

# Function to generate slide content using the OpenAI ChatCompletion API
def generate_presentation_content(topic):
    try:
        messages = [
            {"role": "system", "content": "You are a helpful assistant that generates detailed presentation content."},
            {"role": "user", "content": f"Create a detailed presentation outline for the topic: {topic}. Generate a title, section titles, and for each section, provide 3 numbered points followed by their definitions(explain) in lines as bullet points. Ensure that each point is numbered (1, 2, 3) and the explanations of the points are in bullet format. Create at least 15 slide points."},
             # Additional slide for the concluson
            #{"role": "user", "content": f"Add a slide of Conclusion"},     
            # Additional slide for the future of the topic
            {"role": "user", "content": f"Add a slide titled 'Future of {topic}' with 3 numbered points and their definitions in bullet format."},
        # Q&A slide
            #{"role": "user", "content": f"Add a Q&A slide to encourage audience interaction."},
        # Thank you slide
            #{"role": "user", "content": f"Add a 'Thank You' slide to conclude the presentation."}
        ]

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=1000,
            temperature=0.7
        )
        
        return response['choices'][0]['message']['content'].strip()
    
    except Exception as e:
        print(f"Error generating content: {e}")
        return ""

def generate_subtitle(topic):
    try:
        messages = [
            {"role": "system", "content": "You are a helpful assistant that generates subtitles for presentation slides."},
            {"role": "user", "content": f"Generate a subtitle for the first slide of a presentation on the topic: {topic}."}
        ]

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=100,
            temperature=0.7
        )
        
        return response['choices'][0]['message']['content'].strip()
    
    except Exception as e:
        print(f"Error generating subtitle: {e}")
        return ""

# Function to apply different slide themes based on user choice (using existing templates)
def apply_template(template_choice):
    templates = {
        "Minimalist": "Minimalist.pptx",
        "Corporate": "Corporate.pptx",
        "Creative": "Creative.pptx",
        "Innovative": "Innovative.pptx",
        "Aster": "Aster.pptx"
    }

    template_path = templates.get(template_choice)
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)  # Load the chosen template
    else:
        print(f"Template file {template_path} not found. Falling back to a blank presentation.")
        prs = Presentation()  # Fallback to a blank presentation

    return prs

# Function to add a slide with customized text placement
def add_slide(prs, title_text, content_text, template_choice, image_path=None):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title shape in "Click to add title"
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_shape.text_frame.paragraphs[0].font.bold = True  # Making the title bold
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)

    # Change the title color based on the template choice
    if template_choice == "Creative":
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)  # Yellow color
    elif template_choice == "Aster":
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White color
    else:
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)  # Default navy text

    # Set content box to be wider and centered
    text_left = Inches(1)  # Adjusted left margin for centering
    text_width = Inches(11)  # Increased width
    content_box = slide.shapes.add_textbox(text_left, Inches(2.5), text_width, Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Create a paragraph for the title in the content box and make it bold
    title_paragraph = content_frame.add_paragraph()
    title_paragraph.text = ""  # Set your paragraph title here
    title_paragraph.font.bold = True  # Making the paragraph title bold
    title_paragraph.font.size = Pt(20)  # Adjust the font size of the title
    # title_paragraph.space_after = Pt(12)  # Space after the title
    title_paragraph.alignment = PP_ALIGN.LEFT
    
    # Split content text into bullet points
    bullet_points = content_text.split('\n')  # Assuming each bullet point is on a new line

    # Add each bullet point to the slide with spacing
    for point in bullet_points:
        content_paragraph = content_frame.add_paragraph()
        content_paragraph.text = point.strip()  # Add bullet point text
        
        # Check if the bullet point starts with a digit followed by a period
        if point.strip().split('.')[0].isdigit():
            content_paragraph.font.bold = True  # Apply bold formatting
            
        content_paragraph.font.size = Pt(18)  # Adjust font size for the content text
        content_paragraph.space_after = Pt(10)  # Space after each bullet point
        content_paragraph.alignment = PP_ALIGN.LEFT  # Align to the left

    # Add image if provided
    #if image_path and os.path.exists(image_path):
      #  slide.shapes.add_picture(image_path, Inches(1), Inches(6), width=Inches(8))  # Adjust position and size as needed


 # Add image if provided
    # Add image if provided
    if image_path and os.path.exists(image_path):
    # Image dimensions
        img_width = Inches(5.30)
        img_height = Inches(2.84)

    # Slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

    # Calculate top-left corner to center the image and move it down 4 steps
        step_size = Inches(0.5)  # Adjust the step size as needed (e.g., 0.5 inches)
        top = (slide_height - img_height) / 2 + 3 * step_size

    # Calculate left position to center the image
        left = (slide_width - img_width) / 2

    # Add the image to the slide at the calculated position
        slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)

def create_presentation(topic, slide_count, template_choice):
    prs = apply_template(template_choice)
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    title_slide = prs.slides.add_slide(title_slide_layout)

    # Title text
    title_shape = title_slide.shapes.title
    title_shape.text = topic

    # Generate subtitle using ChatGPT
    subtitle = generate_subtitle(topic)
    subtitle_shape = title_slide.shapes.placeholders[1]  # Placeholder for subtitle
    subtitle_shape.text = subtitle

    # Generate content
    content = generate_presentation_content(topic)
    sections = content.split("\n\n")

    # Limit the sections to the specified slide_count
    sections = sections[:slide_count]

    # Table of Contents slide
    toc_slide_layout = prs.slide_layouts[5]  # Blank layout
    toc_slide = prs.slides.add_slide(toc_slide_layout)

    # Set "Table of Contents" title in the title placeholder
    toc_title_shape = toc_slide.shapes.title
    toc_title_shape.text = "Table of Contents"
    toc_title_shape.text_frame.paragraphs[0].font.size = Pt(32)  # Title font size

    # Create content box below the title with some spacing
    toc_content_box = toc_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))  # Adjusted top margin for spacing
    toc_content_frame = toc_content_box.text_frame
    toc_content_frame.word_wrap = True

    for i, section in enumerate(sections):
        lines = section.strip().split("\n")
        if not lines:
            continue
        
        title = lines[0]
        content_text = "\n".join(lines[1:]).strip()

        # Add the slide with text placement
        add_slide(prs, title, content_text, template_choice)  # Pass template_choice here

        # Add to the table of contents
        p = toc_content_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(18)

    # Add Q&A Slide with image
    add_slide(prs, "Q&A", "Questions & Answers", template_choice, image_path="Q&As.jfif")

    # Add Thank You Slide with image
    add_slide(prs, "Thank You", "Thank You for Your Attention!", template_choice, image_path="Thank You.tiff")


    return prs

# Function to remove specified slides from a presentation
def remove_slides(prs, indices_to_remove):
    indices_to_remove = [0, 3]
    # Sort the indices in descending order to avoid re-indexing issues
    indices_to_remove = sorted(indices_to_remove, reverse=True)

    # Access the internal slide list
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)  # Convert to a list

    # Remove the slides at specified indices
    for idx in indices_to_remove:
        if idx < len(slides):
            xml_slides.remove(slides[idx])  # Remove slide at given index
        else:
            print(f"Slide index {idx} is out of range")

# Save and test the presentation
def save_presentation(prs, filename="presentation.pptx"):
    prs.save(filename)
    print(f"Presentation saved as {filename}")

# Function for voice input
def voice_input(prompt):
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        print(prompt)
        audio = recognizer.listen(source)
        
        try:
            text = recognizer.recognize_google(audio)
            print(f"You said: {text}")
            return text
        except sr.UnknownValueError:
            print("Sorry, I could not understand the audio.")
            return ""
        except sr.RequestError as e:
            print(f"Could not request results from Google Speech Recognition service; {e}")
            return ""

# FastAPI application
app = FastAPI()

@app.post("/generate_presentation/")
async def generate_presentation_endpoint(
    topic: str = Form(...),
    slide_count: int = Form(...),
    template_choice: str = Form(...)
):
    """AI Presentation Generator.
            Choose a template:
        Minimalist,
        Corporate,
        Creative,
        Innovative,
        Aster"""
    if slide_count < 5 or slide_count > 18:
        raise HTTPException(status_code=400, detail="Slide count must be between 5 and 18.")

    try:
        # Generate the presentation
        prs = create_presentation(topic, slide_count, template_choice)

        # Remove slides 0 and 3
        remove_slides(prs, [0, 3])

        # Save the presentation
        filename = f"{topic.replace(' ', '_')}.pptx"
        save_presentation(prs, filename)

        # Return the presentation file for download
        return FileResponse(
            path=filename,
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
@app.get("/")
def root():
    return {"message": "Welcome to the Presentation Generator API"}
